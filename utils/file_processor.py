import os
import fitz
import logging
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import requests
from bs4 import BeautifulSoup
from urllib.parse import urlparse
import tempfile
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import re
from io import BytesIO
from PyPDF2 import PdfReader
import markdown
from weasyprint import HTML
from weasyprint.text.fonts import FontConfiguration
from docx import Document
from striprtf.striprtf import rtf_to_text
import zipfile
import xml.etree.ElementTree as ET
import magic

logging.basicConfig(level=logging.ERROR)
logger = logging.getLogger(__name__)


def process_file(input_data):
    logger.debug(f"Starting to process input: {input_data}")
    try:
        file_type = magic.from_file(input_data, mime=True)
        file_extension = os.path.splitext(input_data)[1].lower()
        logger.debug(
            f"Detected file type: {file_type}, File extension: {file_extension}"
        )

        if file_type in [
                'application/octet-stream', 'text/plain', 'application/zip'
        ]:
            if file_extension == '.pptx':
                file_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            elif file_extension == '.key':
                file_type = 'application/x-iwork-keynote-sffkey'
            elif file_extension == '.md':
                file_type = 'text/markdown'

        temp_pdf_path = None
        video_tracks = []
        audio_tracks = []
        original_type = file_type

        if file_type == 'application/pdf':
            temp_pdf_path = input_data
        elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            temp_pdf_path, video_tracks, audio_tracks = convert_to_pdf(
                input_data)
        elif file_type == 'text/markdown' or file_extension == '.md':
            result = convert_markdown_to_pdf(input_data)
            # Use the result directly instead of further processing
            return result
        elif file_type == 'application/x-iwork-keynote-sffkey' or file_extension == '.key':
            temp_pdf_path = convert_keynote_to_pdf(input_data)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

        result = process_pdf(temp_pdf_path)

        if result:
            result.update({
                'original_type': original_type,
                'type': 'application/pdf',
                'temp_file_path': temp_pdf_path,
                'video_tracks': video_tracks,
                'audio_tracks': audio_tracks
            })
            return result
        else:
            raise ValueError("Processing failed to produce a result")
    except Exception as e:
        logger.error(f"Error in process_file: {str(e)}", exc_info=True)
        return {'error': str(e), 'type': 'unknown'}


def process_pdf(pdf_path):
    try:
        doc = fitz.open(pdf_path)
        num_pages = len(doc)
        content = []
        fonts = set()

        for page in doc:
            content.append(page.get_text())

            # Extract fonts from each page
            for font in page.get_fonts():
                # font[3] is the font name
                fonts.add(font[3])

        doc.close()

        return {
            'type': 'pdf',
            'num_slides': num_pages,
            'content': content,
            'fonts': list(fonts)
        }
    except Exception as e:
        logger.error(f"Error processing PDF: {str(e)}", exc_info=True)
        return {'error': f"Failed to process PDF: {str(e)}"}


def convert_to_pdf(input_file):
    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
        output_file = temp_pdf.name

    try:
        prs = Presentation(input_file)
        pdf = canvas.Canvas(output_file, pagesize=letter)

        video_tracks = []
        audio_tracks = []

        for slide_index, slide in enumerate(prs.slides):
            pdf.setFont("Helvetica", 12)
            pdf.drawString(100, 700, f"Slide {slide_index + 1}")
            y_position = 680
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        text = paragraph.text
                        pdf.drawString(100, y_position,
                                       text[:80])  # Limit text to 80 chars
                        y_position -= 20
                elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                    if hasattr(shape, 'media_format') and hasattr(
                            shape.media_format, 'mime_type'):
                        mime_type = shape.media_format.mime_type
                        if mime_type.startswith('video'):
                            video_tracks.append(
                                f"Video on slide {slide_index + 1}")
                            pdf.drawString(100, y_position, "[Video]")
                        elif mime_type.startswith('audio'):
                            audio_tracks.append(
                                f"Audio on slide {slide_index + 1}")
                            pdf.drawString(100, y_position, "[Audio]")
                    y_position -= 20

                if y_position < 50:
                    pdf.showPage()
                    y_position = 700
            pdf.showPage()
        pdf.save()
        logger.debug(f"Successfully converted {input_file} to PDF")
        return output_file, video_tracks, audio_tracks
    except Exception as e:
        logger.error(f"Error converting {input_file} to PDF: {str(e)}",
                     exc_info=True)
        raise


def convert_keynote_to_pdf(input_file):
    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
        output_file = temp_pdf.name

    try:
        slides = extract_text_from_keynote(input_file)

        pdf = canvas.Canvas(output_file, pagesize=letter)
        pdf.setFont("Helvetica", 12)

        for i, slide_content in enumerate(slides):
            pdf.drawString(100, 750, f"Slide {i + 1}")
            y = 720
            for line in slide_content.split('\n'):
                pdf.drawString(100, y, line)
                y -= 20
                if y < 50:
                    pdf.showPage()
                    y = 750
            pdf.showPage()

        pdf.save()
        return output_file
    except Exception as e:
        logger.error(f"Error converting Keynote to PDF: {str(e)}",
                     exc_info=True)
        raise


def extract_text_from_keynote(keynote_file):
    slides = []
    with zipfile.ZipFile(keynote_file, 'r') as zip_ref:
        for filename in zip_ref.namelist():
            if filename.startswith('Data/Slide') and filename.endswith(
                    '.apxl'):
                with zip_ref.open(filename) as file:
                    tree = ET.parse(file)
                    root = tree.getroot()
                    slide_content = ""
                    for text_elem in root.iter(
                            '{http://developer.apple.com/namespaces/keynote2}text'
                    ):
                        slide_content += text_elem.text + "\n" if text_elem.text else ""
                    slides.append(slide_content)
    return slides


def convert_markdown_to_pdf(markdown_path):
    try:
        with open(markdown_path, 'r', encoding='utf-8') as md_file:
            md_content = md_file.read()
        html = markdown.markdown(md_content)
        with tempfile.NamedTemporaryFile(suffix='.pdf',
                                         delete=False) as temp_pdf:
            temp_pdf_path = temp_pdf.name
        font_config = FontConfiguration()
        HTML(string=html).write_pdf(temp_pdf_path, font_config=font_config)

        # Process the PDF to get additional information
        pdf_info = process_pdf(temp_pdf_path)

        return {
            'type': 'application/pdf',
            'original_type': 'markdown',
            'num_slides': pdf_info['num_slides'],
            'content': pdf_info['content'],
            'fonts': pdf_info.get('fonts', []),
            'temp_file_path': temp_pdf_path
        }
    except Exception as e:
        logger.error(f"Error converting Markdown to PDF: {str(e)}",
                     exc_info=True)
        raise


def process_url(url):
    try:
        parsed_url = urlparse(url)
        if 'docs.google.com' in parsed_url.netloc and 'presentation' in parsed_url.path:
            return process_google_slides(url)
        elif 'figma.com' in parsed_url.netloc:
            return process_figma(url)
        elif 'canva.com' in parsed_url.netloc:
            return process_canva(url)
        else:
            raise ValueError("Unsupported URL type")
    except Exception as e:
        logger.error(f"Error processing URL: {str(e)}", exc_info=True)
        return {'error': str(e), 'type': 'unknown', 'url': url}


def process_figma(url):
    try:
        response = requests.get(url)
        response.raise_for_status()

        # Save as PDF
        with tempfile.NamedTemporaryFile(suffix='.pdf',
                                         delete=False) as temp_pdf:
            temp_pdf_path = temp_pdf.name

        HTML(string=response.text).write_pdf(temp_pdf_path)

        result = process_pdf(temp_pdf_path)
        result.update({
            'original_type': 'figma',
            'type': 'application/pdf',
            'url': url,
            'temp_file_path': temp_pdf_path
        })
        return result
    except Exception as e:
        logger.error(f"Error processing Figma URL: {str(e)}", exc_info=True)
        return {'error': str(e), 'type': 'figma'}


def process_canva(url):
    try:
        response = requests.get(url)
        response.raise_for_status()

        # Save as PDF
        with tempfile.NamedTemporaryFile(suffix='.pdf',
                                         delete=False) as temp_pdf:
            temp_pdf_path = temp_pdf.name

        HTML(string=response.text).write_pdf(temp_pdf_path)

        result = process_pdf(temp_pdf_path)
        result.update({
            'original_type': 'canva',
            'type': 'application/pdf',
            'url': url,
            'temp_file_path': temp_pdf_path
        })
        return result
    except Exception as e:
        logger.error(f"Error processing Canva URL: {str(e)}", exc_info=True)
        return {'error': str(e), 'type': 'canva'}


def process_google_slides(url):
    try:
        # Extract presentation ID from URL
        match = re.search('/d/([a-zA-Z0-9-_]+)', url)
        if not match:
            raise ValueError("Invalid Google Slides URL")
        presentation_id = match.group(1)

        # Construct the export URL
        export_url = f"https://docs.google.com/presentation/d/{presentation_id}/export/pdf"

        # Download the PDF
        response = requests.get(export_url)
        if response.status_code != 200:
            raise Exception(
                "Failed to download the presentation. Make sure it's public and the URL is correct."
            )

        # Save the PDF content to a temporary file
        with tempfile.NamedTemporaryFile(suffix='.pdf',
                                         delete=False) as temp_pdf:
            temp_pdf_path = temp_pdf.name
            temp_pdf.write(response.content)

        result = process_pdf(temp_pdf_path)
        result.update({
            'original_type': 'google_slides',
            'type': 'application/pdf',
            'url': url,
            'temp_file_path': temp_pdf_path
        })
        return result
    except Exception as e:
        logger.error(f"Error processing Google Slides: {str(e)}",
                     exc_info=True)
        return {'error': str(e), 'type': 'google_slides', 'url': url}


# Main execution
if __name__ == "__main__":
    # You can add test code here to run the script directly
    pass
