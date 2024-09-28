import os
from utils.file_processor import process_file
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from odf.opendocument import OpenDocumentPresentation
from odf.style import Style, TextProperties, ParagraphProperties
from odf.draw import Page, Frame, TextBox
from odf.text import P

def create_sample_pptx():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Sample PPTX Presentation"
    subtitle.text = "This is a sample PPTX file for testing purposes."
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_pptx:
        prs.save(temp_pptx.name)
    return temp_pptx.name

def create_sample_otp():
    doc = OpenDocumentPresentation()
    
    # Create a style for the text
    textstyle = Style(name="Standard", family="paragraph")
    textstyle.addElement(TextProperties(fontsize="24pt", fontweight="bold"))
    doc.styles.addElement(textstyle)

    # Create a page
    page = Page()
    
    # Add a text box to the page
    frame = Frame(width="500pt", height="100pt", x="50pt", y="50pt")
    textbox = TextBox()
    
    # Add text to the text box
    p = P(stylename=textstyle)
    p.addText("Sample OTP Presentation")
    textbox.addElement(p)
    
    p = P(stylename=textstyle)
    p.addText("This is a sample OTP file for testing purposes.")
    textbox.addElement(p)
    
    frame.addElement(textbox)
    page.addElement(frame)
    doc.body.addElement(page)

    with tempfile.NamedTemporaryFile(suffix='.otp', delete=False) as temp_otp:
        doc.save(temp_otp.name)
    return temp_otp.name

def test_pptx_processing():
    pptx_file = create_sample_pptx()
    result = process_file(pptx_file)
    print("PPTX Processing Result:")
    print(f"Number of slides: {result['num_slides']}")
    print(f"Content: {result['content']}")
    os.unlink(pptx_file)
    os.unlink(result['temp_file_path'])

def test_otp_processing():
    otp_file = create_sample_otp()
    result = process_file(otp_file)
    print("OTP Processing Result:")
    print(f"Number of slides: {result['num_slides']}")
    print(f"Content: {result['content']}")
    os.unlink(otp_file)
    os.unlink(result['temp_file_path'])

if __name__ == "__main__":
    test_pptx_processing()
    test_otp_processing()
